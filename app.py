import sys
import re
import json
import html
import time
import queue
import logging

logger = logging.getLogger(__name__)
import platform
from pathlib import Path

import cv2
import numpy as np

from PySide6.QtCore import Qt, QTimer, QThread, Signal
from PySide6.QtGui import (
    QImage,
    QPixmap,
    QPainter,
    QPen,
    QBrush,
    QColor,
    QFont,
    QLinearGradient,
)
from PySide6.QtWidgets import (
    QApplication,
    QMainWindow,
    QWidget,
    QLabel,
    QPushButton,
    QTextEdit,
    QLineEdit,
    QHBoxLayout,
    QVBoxLayout,
    QSplitter,
    QTabWidget,
    QProgressBar,
    QStackedWidget,
    QFileDialog,
    QFrame,
    QGridLayout,
)

# Import new modules
from llm_brain import LLMBrain
from tts_engine import TTSEngine
from interrupt_handler import InterruptHandler
from video_avatar import VideoAvatar
from language_manager import LanguageManager
from conversation_engine import ConversationEngine
from voice_listener import ContinuousVoiceListener
from gesture_detector import GestureDetector
from analytics_manager import AnalyticsManager



APP_DIR = Path(__file__).resolve().parent
PPT_PATH = APP_DIR / "default_presentation.pptx"
VOSK_MODEL_PATH = APP_DIR / "models" / "vosk-model-small-en-us-0.15"

class PPTManager:
    def __init__(self, ppt_path: Path):
        self.ppt_path = ppt_path
        self.slides: list[dict] = []
        self.ensure_sample_ppt()
        self.load()

    def ensure_sample_ppt(self):
        if self.ppt_path.exists():
            return

        from pptx import Presentation
        prs = Presentation()

        data = [
            (
                "Welcome to the Intelligence Lab",
                "This demo shows a camera based interactive anime agent. "
                "When a person is detected, the agent greets the visitor and starts the presentation automatically.",
            ),
            (
                "What the Agent Can Do",
                "The agent can detect a visitor using webcam input, display an animated avatar, "
                "read slide content aloud, accept typed questions, and answer using the presentation content.",
            ),
            (
                "Core Technologies",
                "The demo uses PySide6 for the user interface, OpenCV for webcam face detection, "
                "python-pptx for reading PowerPoint text, Windows speech or pyttsx3 for speech output, "
                "and Vosk for optional offline voice questions.",
            ),
            (
                "Demo Flow",
                "Step one: camera starts. Step two: visitor is detected. Step three: anime agent appears and greets the user. "
                "Step four: presentation starts. Step five: user asks questions by chat or microphone.",
            ),
            (
                "Possible Enhancements",
                "The demo can be extended with company branding, real slide image rendering, local LLM answers, "
                "richer 2D avatar animation, hand gesture control, and visitor analytics.",
            ),
            (
                "Thank You",
                "This completes the default demo. I am ready to answer questions from the presentation.",
            ),
        ]

        for title, body in data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = body

        prs.save(self.ppt_path)

    def load(self):
        from pptx import Presentation
        prs = Presentation(str(self.ppt_path))
        self.slides.clear()

        for index, slide in enumerate(prs.slides, start=1):
            parts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    parts.append(shape.text.strip())

            text = "\n".join(parts).strip()

            if not text:
                text = "This slide has no readable text. It may contain only images or diagrams."

            title = parts[0].splitlines()[0] if parts else f"Slide {index}"

            self.slides.append(
                {
                    "number": index,
                    "title": title,
                    "text": text,
                }
            )

    def load_presentation(self, new_path: str):
        self.ppt_path = Path(new_path)
        self.load()




class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()

        self.setWindowTitle("Intelligence Lab Demo Agent - Human-Like AI Presenter")
        self.resize(900, 600)

        # 1. Initialize core non-UI managers
        self.ppt = PPTManager(PPT_PATH)
        self.llm = LLMBrain()
        self.lm = LanguageManager()
        self.analytics = AnalyticsManager()
        
        # 2. Initialize TTS thread and signals
        self.tts = TTSEngine()
        self.tts.status.connect(self.set_status)
        self.tts.start()

        self.interrupt_handler = InterruptHandler(self.tts)
        
        # 3. Setup conversation flow engine
        self.conv_engine = ConversationEngine(
            llm_brain=self.llm,
            tts_engine=self.tts,
            ppt_manager=self.ppt,
            interrupt_handler=self.interrupt_handler,
            language_manager=self.lm
        )

        self.demo_started = False
        self.speaker_muted = False
        self.frame_count = 0
        self.face_stable_count = 0
        self.no_face_frames = 0
        self.last_detection_label = "None"
        self.camera_index = None

        self.frontal_cascade = self.load_cascade("haarcascade_frontalface_default.xml")
        self.profile_cascade = self.load_cascade("haarcascade_profileface.xml")
        self.upperbody_cascade = self.load_cascade("haarcascade_upperbody.xml")

        self.cap = None

        # 4. Initialize Continuous Voice Activity Listener and Hand Gesture Detector
        self.voice_listener = None
        self.gesture_detector = GestureDetector()
        

        QTimer.singleShot(500, self.setup_voice_listener)

        # 5. Build components and UI
        self.build_ui()
        QTimer.singleShot(100, self.setup_camera)
        
        # Connect conversation engine events to UI
        self.conv_engine.state_changed.connect(self.on_state_changed)
        self.conv_engine.text_to_speak.connect(self.on_text_to_speak)
        self.conv_engine.chat_message.connect(self.add_chat)
        self.conv_engine.slide_navigation.connect(self.show_slide)
        self.conv_engine.ui_status_update.connect(self.set_status)
        
        # Connect TTS speaking state for lip sync and echo prevention
        self.tts.speaking_state_changed.connect(self.on_speaking_state_changed)

        self.show_slide(0)

        # Check active LLM backend and print setup tip
        if self.conv_engine.llm.gemini_key:
            self.set_status("Active: Google Gemini API (Cloud). Responses will be instant!")
        elif self.conv_engine.llm.openai_key:
            self.set_status("Active: OpenAI API (Cloud). Responses will be instant!")
        else:
            self.set_status("Active: Local Ollama. Tip: Set GEMINI_API_KEY in terminal for <1s responses.")

    def load_cascade(self, filename: str):
        path = cv2.data.haarcascades + filename
        cascade = cv2.CascadeClassifier(path)
        if cascade.empty():
            return None
        return cascade

    def setup_voice_listener(self):
        if not VOSK_MODEL_PATH.exists():
            self.set_status("Vosk model not found. Continuous voice input disabled.")
            return

        self.voice_listener = ContinuousVoiceListener(str(VOSK_MODEL_PATH))
        self.voice_listener.partial_speech.connect(self.on_partial_voice_text)
        self.voice_listener.speech_completed.connect(self.on_voice_text)
        self.voice_listener.status.connect(self.set_status)
        self.voice_listener.error.connect(self.set_status)
        self.voice_listener.start()
        # Sync initial mic state with UI button toggle
        self.voice_listener.set_enabled(self.mic_toggle.isChecked())

    def build_ui(self):
        root = QWidget()
        self.setCentralWidget(root)
        
        layout = QVBoxLayout(root)
        layout.setContentsMargins(10, 8, 10, 8)
        layout.setSpacing(8)
        
        # Header Row (Title + Nav Toggle Buttons)
        header = QHBoxLayout()
        
        title = QLabel("🤖 Intelligence Lab Interactive Presenter - Aiko")
        title.setObjectName("Title")
        title.setFont(QFont("Segoe UI", 16, QFont.Bold))
        title.setStyleSheet("color: #ffffff; padding: 0;")
        header.addWidget(title)
        
        header.addStretch()
        
        # Navigation toggle buttons (Pill style)
        self.btn_nav_kiosk = QPushButton("👥 Live Kiosk")
        self.btn_nav_kiosk.setObjectName("NavBtnActive")
        self.btn_nav_kiosk.setStyleSheet("background: #4759e5; color: white; padding: 8px 18px; border-radius: 18px; font-weight: bold; border: none;")
        self.btn_nav_kiosk.clicked.connect(self.show_kiosk_view)
        
        self.btn_nav_dash = QPushButton("📊 Analytics Dashboard")
        self.btn_nav_dash.setObjectName("NavBtnInactive")
        self.btn_nav_dash.setStyleSheet("background: #1a1f3c; color: #a4b2fc; padding: 8px 18px; border-radius: 18px; font-weight: bold; border: 1px solid #2e3557;")
        self.btn_nav_dash.clicked.connect(self.show_dashboard_view)
        
        header.addWidget(self.btn_nav_kiosk)
        header.addWidget(self.btn_nav_dash)
        
        layout.addLayout(header)
        
        # Stacked Widget
        self.stacked_widget = QStackedWidget()
        layout.addWidget(self.stacked_widget, stretch=1)
        
        # Status Bar at bottom
        self.status = QLabel("Status: starting...")
        self.status.setObjectName("Status")
        layout.addWidget(self.status)
        
        # Build and add sub-views
        kiosk_view = self.build_kiosk_tab()
        self.stacked_widget.addWidget(kiosk_view)
        
        dashboard_view = self.build_dashboard()
        self.stacked_widget.addWidget(dashboard_view)
        
        self.apply_stylesheet()

    def show_kiosk_view(self):
        self.stacked_widget.setCurrentIndex(0)
        self.btn_nav_kiosk.setStyleSheet("background: #4759e5; color: white; padding: 8px 18px; border-radius: 18px; font-weight: bold; border: none;")
        self.btn_nav_dash.setStyleSheet("background: #1a1f3c; color: #a4b2fc; padding: 8px 18px; border-radius: 18px; font-weight: bold; border: 1px solid #2e3557;")

    def show_dashboard_view(self):
        self.stacked_widget.setCurrentIndex(1)
        self.btn_nav_kiosk.setStyleSheet("background: #1a1f3c; color: #a4b2fc; padding: 8px 18px; border-radius: 18px; font-weight: bold; border: 1px solid #2e3557;")
        self.btn_nav_dash.setStyleSheet("background: #4759e5; color: white; padding: 8px 18px; border-radius: 18px; font-weight: bold; border: none;")
        self.refresh_dashboard()

    def build_kiosk_tab(self) -> QWidget:
        kiosk_widget = QWidget()
        outer = QVBoxLayout(kiosk_widget)
        outer.setContentsMargins(0, 0, 0, 0)

        splitter = QSplitter(Qt.Horizontal)
        outer.addWidget(splitter, stretch=1)

        self.video_label = QLabel("Camera starting...")
        self.video_label.setAlignment(Qt.AlignCenter)
        self.video_label.setMinimumWidth(360)
        self.video_label.setObjectName("Panel")
        splitter.addWidget(self.video_label)

        center = QWidget()
        center_layout = QVBoxLayout(center)

        self.slide_title = QLabel("Slide Title")
        self.slide_title.setObjectName("SlideTitle")
        self.slide_title.setWordWrap(True)

        self.slide_text = QTextEdit()
        self.slide_text.setReadOnly(True)
        self.slide_text.setObjectName("SlideText")

        center_layout.addWidget(self.slide_title)
        center_layout.addWidget(self.slide_text, stretch=1)

        # Create a container frame for all controls (Language, Voice, Action buttons)
        control_deck = QFrame()
        control_deck.setObjectName("ControlDeck")
        control_deck_layout = QVBoxLayout(control_deck)
        control_deck_layout.setContentsMargins(12, 10, 12, 10)
        control_deck_layout.setSpacing(10)

        # Row 1: Language selection Group
        lang_row = QHBoxLayout()
        lang_title = QLabel("Language:")
        lang_title.setFont(QFont("Segoe UI", 10, QFont.Bold))
        lang_row.addWidget(lang_title)

        self.lang_buttons = []
        for lang_info in self.lm.get_supported_languages():
            btn = QPushButton(f"{lang_info['flag_emoji']} {lang_info['native_name']}")
            btn.setObjectName("LangBtn")
            btn.setCheckable(True)
            if lang_info['code'] == self.lm.get_language():
                btn.setChecked(True)
            btn.clicked.connect(lambda checked=False, code=lang_info['code']: self.select_language(code))
            lang_row.addWidget(btn)
            self.lang_buttons.append((btn, lang_info['code']))
        
        control_deck_layout.addLayout(lang_row)

        # Row 2: Settings (Mic Toggle + Load PPTX side-by-side)
        settings_row = QHBoxLayout()
        
        self.mic_toggle = QPushButton("🎙️ Voice Input: OFF")
        self.mic_toggle.setCheckable(True)
        self.mic_toggle.setChecked(False)
        self.mic_toggle.setObjectName("AudioBtn")
        self.mic_toggle.clicked.connect(self.toggle_mic)
        settings_row.addWidget(self.mic_toggle, stretch=1)
        
        self.upload_btn = QPushButton("📂 Load Custom PPTX")
        self.upload_btn.setObjectName("UploadBtn")
        self.upload_btn.clicked.connect(self.upload_presentation)
        settings_row.addWidget(self.upload_btn, stretch=1)

        control_deck_layout.addLayout(settings_row)

        # Row 3: Action Buttons (2x2 Grid Layout)
        grid_layout = QGridLayout()
        grid_layout.setSpacing(6)

        self.start_btn = QPushButton("🚀 Start Presenter")
        self.start_btn.setObjectName("StartBtn")
        self.repeat_btn = QPushButton("🔁 Repeat Slide")
        self.next_btn = QPushButton("➡ Next Slide")
        self.restart_btn = QPushButton("▶ Restart Presentation")

        self.start_btn.clicked.connect(self.on_start_or_resume_clicked)
        self.repeat_btn.clicked.connect(self.repeat_slide)
        self.next_btn.clicked.connect(self.next_slide_manual)
        self.restart_btn.clicked.connect(self.restart_presentation)

        grid_layout.addWidget(self.start_btn, 0, 0)
        grid_layout.addWidget(self.repeat_btn, 0, 1)
        grid_layout.addWidget(self.next_btn, 1, 0)
        grid_layout.addWidget(self.restart_btn, 1, 1)

        control_deck_layout.addLayout(grid_layout)
        
        center_layout.addWidget(control_deck)
        splitter.addWidget(center)

        # Use VideoAvatar for realistic human AI presenter
        self.avatar = VideoAvatar()
        splitter.addWidget(self.avatar)

        splitter.setSizes([80, 420, 400])

        self.chat = QTextEdit()
        self.chat.setReadOnly(True)
        self.chat.setMaximumHeight(110)
        self.chat.setObjectName("ChatBox")

        outer.addWidget(self.chat)

        input_row = QHBoxLayout()

        self.input_box = QLineEdit()
        self.input_box.setPlaceholderText("Ask a question about the presentation...")
        self.input_box.setFixedHeight(36)
        self.input_box.setObjectName("ChatInputBox")
        self.input_box.setStyleSheet("""
            QLineEdit#ChatInputBox {
                background: #ffffff;
                border: 1px solid #3c457d;
                border-radius: 6px;
                color: #1a1f36;
                font-size: 13px;
                padding-left: 10px;
            }
        """)

        self.send_btn = QPushButton("Send")
        self.send_btn.setFixedHeight(36)
        
        self.input_box.returnPressed.connect(self.handle_typed_question)
        self.send_btn.clicked.connect(self.handle_typed_question)

        input_row.addWidget(self.input_box, stretch=1)
        input_row.addWidget(self.send_btn)
        outer.addLayout(input_row)
        return kiosk_widget

    def build_dashboard(self) -> QWidget:
        dashboard_widget = QWidget()
        layout = QVBoxLayout(dashboard_widget)
        layout.setContentsMargins(15, 15, 15, 15)
        
        # Title
        title = QLabel("📊 Enterprise Kiosk Analytics Dashboard")
        title.setFont(QFont("Segoe UI", 16, QFont.Bold))
        title.setStyleSheet("color: #ffffff; padding-bottom: 10px;")
        layout.addWidget(title)
        
        # Row 1: Stats Cards (Visitors, Questions, Avg Duration)
        stats_layout = QHBoxLayout()
        stats_layout.setSpacing(15)
        
        # Card 1: Visitors
        self.card_visitors = QWidget()
        self.card_visitors.setObjectName("StatsCard1")
        cv_layout = QVBoxLayout(self.card_visitors)
        lbl_v = QLabel("Total Visitors")
        lbl_v.setStyleSheet("font-size: 13px; color: #a4b2fc; font-weight: 500;")
        self.val_visitors = QLabel("0")
        self.val_visitors.setStyleSheet("font-size: 32px; font-weight: bold; color: #ffffff; margin-top: 5px;")
        cv_layout.addWidget(lbl_v)
        cv_layout.addWidget(self.val_visitors)
        stats_layout.addWidget(self.card_visitors)
        
        # Card 2: Questions
        self.card_questions = QWidget()
        self.card_questions.setObjectName("StatsCard2")
        cq_layout = QVBoxLayout(self.card_questions)
        lbl_q = QLabel("Questions Answered")
        lbl_q.setStyleSheet("font-size: 13px; color: #c084fc; font-weight: 500;")
        self.val_questions = QLabel("0")
        self.val_questions.setStyleSheet("font-size: 32px; font-weight: bold; color: #ffffff; margin-top: 5px;")
        cq_layout.addWidget(lbl_q)
        cq_layout.addWidget(self.val_questions)
        stats_layout.addWidget(self.card_questions)
        
        # Card 3: Avg Duration
        self.card_duration = QWidget()
        self.card_duration.setObjectName("StatsCard3")
        cd_layout = QVBoxLayout(self.card_duration)
        lbl_d = QLabel("Avg. Visitor Engagement")
        lbl_d.setStyleSheet("font-size: 13px; color: #10b981; font-weight: 500;")
        self.val_duration = QLabel("0.0s")
        self.val_duration.setStyleSheet("font-size: 32px; font-weight: bold; color: #ffffff; margin-top: 5px;")
        cd_layout.addWidget(lbl_d)
        cd_layout.addWidget(self.val_duration)
        stats_layout.addWidget(self.card_duration)
        
        layout.addLayout(stats_layout)
        
        # Row 2: Splitter for charts and logs
        splitter = QSplitter(Qt.Horizontal)
        splitter.setContentsMargins(0, 10, 0, 0)
        layout.addWidget(splitter, stretch=1)
        
        # Left Panel: Slide Popularity
        left_panel = QWidget()
        left_panel.setObjectName("Panel")
        left_layout = QVBoxLayout(left_panel)
        left_title = QLabel("🔥 Slide View Popularity")
        left_title.setFont(QFont("Segoe UI", 12, QFont.Bold))
        left_title.setStyleSheet("color: #ffffff; padding-bottom: 8px;")
        left_layout.addWidget(left_title)
        
        self.slide_list_widget = QWidget()
        self.slide_list_layout = QVBoxLayout(self.slide_list_widget)
        self.slide_list_layout.setContentsMargins(0, 0, 0, 0)
        self.slide_list_layout.addStretch()
        left_layout.addWidget(self.slide_list_widget, stretch=1)
        
        splitter.addWidget(left_panel)
        
        # Right Panel: Recent Questions Log
        right_panel = QWidget()
        right_panel.setObjectName("Panel")
        right_layout = QVBoxLayout(right_panel)
        right_title = QLabel("💬 Recent Visitor Questions (Live)")
        right_title.setFont(QFont("Segoe UI", 12, QFont.Bold))
        right_title.setStyleSheet("color: #ffffff; padding-bottom: 8px;")
        right_layout.addWidget(right_title)
        
        self.recent_questions_text = QTextEdit()
        self.recent_questions_text.setReadOnly(True)
        self.recent_questions_text.setStyleSheet(
            "background: #080a12; border: 1px solid #1e293b; border-radius: 8px; font-size: 13px; color: #cbd5e1; padding: 8px;"
        )
        right_layout.addWidget(self.recent_questions_text, stretch=1)
        
        splitter.addWidget(right_panel)
        splitter.setSizes([600, 600])
        
        return dashboard_widget



    def refresh_dashboard(self):
        stats = self.analytics.get_stats()
        
        # Update metrics
        self.val_visitors.setText(str(stats["total_visitors"]))
        self.val_questions.setText(str(stats["total_questions"]))
        self.val_duration.setText(f"{stats['avg_session_time']}s")
        
        # Refresh Slide views horizontal bar chart
        # Clear the slide list layout
        while self.slide_list_layout.count() > 0:
            item = self.slide_list_layout.takeAt(0)
            widget = item.widget()
            if widget:
                widget.deleteLater()
                
        # Populate progress bars
        slide_views = stats["slide_views"]
        max_views = max([x[1] for x in slide_views]) if slide_views else 1
        
        if not slide_views:
            lbl = QLabel("No slide view statistics recorded yet.")
            lbl.setAlignment(Qt.AlignCenter)
            lbl.setStyleSheet("color: #64748b; font-size: 14px; margin-top: 40px;")
            self.slide_list_layout.addWidget(lbl)
        else:
            for title, count in slide_views[:6]: # top 6 slides
                row = QWidget()
                row_layout = QHBoxLayout(row)
                row_layout.setContentsMargins(0, 4, 0, 4)
                
                title_lbl = QLabel(title)
                title_lbl.setWordWrap(True)
                title_lbl.setMaximumWidth(220)
                title_lbl.setStyleSheet("color: #eef2ff; font-weight: 500; font-size: 13px;")
                
                bar = QProgressBar()
                bar.setMinimum(0)
                bar.setMaximum(max_views)
                bar.setValue(count)
                
                count_lbl = QLabel(f"{count} views")
                count_lbl.setStyleSheet("color: #a4b2fc; font-weight: bold; min-width: 60px; font-size: 13px;")
                count_lbl.setAlignment(Qt.AlignRight | Qt.AlignVCenter)
                
                row_layout.addWidget(title_lbl)
                row_layout.addWidget(bar, stretch=1)
                row_layout.addWidget(count_lbl)
                
                self.slide_list_layout.addWidget(row)
                
        self.slide_list_layout.addStretch()
        
        # Refresh recent questions text
        questions = stats["recent_questions"]
        if not questions:
            self.recent_questions_text.setHtml("<p style='color: #64748b; text-align: center; margin-top: 40px;'>No questions logged yet.</p>")
        else:
            html_text = ""
            for q in questions:
                # Color code language label
                lang_color = "#3b82f6" if q["language"].lower() == "english" else "#10b981"
                html_text += (
                    f"<div style='margin-bottom: 12px; border-bottom: 1px solid #1e293b; padding-bottom: 8px;'>"
                    f"  <span style='color: #94a3b8; font-size: 11px;'>[{q['timestamp']}]</span> "
                    f"  <span style='background-color: {lang_color}; color: #ffffff; font-size: 11px; font-weight: bold; border-radius: 4px; padding: 2px 6px;'>{q['language'].upper()}</span>"
                    f"  <p style='margin: 4px 0 0 0; font-weight: 500; color: #e2e8f0; font-size: 13px;'>{html.escape(q['question'])}</p>"
                    f"</div>"
                )
            self.recent_questions_text.setHtml(html_text)

    def apply_stylesheet(self):
        self.setStyleSheet(
            """
            QWidget {
                background: #111424;
                color: #eef2ff;
                font-family: Segoe UI, Arial;
                font-size: 13px;
            }

            QLabel#Title {
                font-size: 20px;
                font-weight: 700;
                padding: 10px;
                color: #ffffff;
            }

            QLabel#Panel {
                background: #080a12;
                border: 1px solid #2e3557;
                border-radius: 12px;
                padding: 10px;
            }

            QLabel#SlideTitle {
                font-size: 22px;
                font-weight: 700;
                color: #ffffff;
                padding: 10px;
                background: #1d2345;
                border-radius: 12px;
            }

            QTextEdit#SlideText {
                background: #f7f9fc;
                color: #1a1f36;
                border-radius: 12px;
                padding: 14px;
                font-size: 18px;
                font-weight: 500;
            }

            QTextEdit#ChatBox {
                background: #171c33;
                border: 1px solid #2e3557;
                border-radius: 10px;
                padding: 8px;
                font-size: 14px;
            }

            QLineEdit {
                background: #ffffff;
                color: #1a1f36;
                border-radius: 8px;
                padding: 10px;
                font-size: 15px;
            }

            QFrame#ControlDeck {
                background: #0b0d19;
                border: 1px solid #1e2547;
                border-radius: 12px;
            }

            QPushButton {
                background: #181d3d;
                color: #a4b2fc;
                border: 1px solid #2e3557;
                border-radius: 6px;
                padding: 10px 14px;
                font-weight: 700;
            }

            QPushButton:hover {
                background: #252b54;
                border-color: #3b4585;
                color: #ffffff;
            }

            QPushButton#StartBtn {
                background: #10b981;
                border: 1px solid #059669;
                color: #ffffff;
            }

            QPushButton#StartBtn:hover {
                background: #059669;
            }

            QPushButton#UploadBtn {
                background: #4759e5;
                border: 1px solid #3c4cd4;
                color: #ffffff;
            }

            QPushButton#UploadBtn:hover {
                background: #3646c0;
            }
            
            QPushButton#LangBtn {
                background: #151a35;
                border: 1px solid #2e3557;
                padding: 6px 12px;
                font-size: 12px;
            }
            
            QPushButton#LangBtn:hover {
                background: #252b54;
                border-color: #3b4585;
            }

            QPushButton#LangBtn:checked {
                background: #4759e5;
                border: 1px solid #3c4cd4;
                color: #ffffff;
            }

            QPushButton#AudioBtn {
                background: #22c55e;
                border: 1px solid #16a34a;
                padding: 6px 12px;
                font-size: 12px;
            }
            
            QPushButton#AudioBtn:hover {
                background: #252b54;
                border-color: #3b4585;
            }

            QPushButton#AudioBtn:checked {
                background: #10b981;
                border: 1px solid #059669;
                color: #ffffff;
            }
            
            QPushButton#AudioBtn:!checked {
                background: #1e1b4b;
                border: 1px solid #312e81;
                color: #cbd5e1;
            }

            QLabel#Status {
                color: #a4b2fc;
                padding: 6px;
            }

            QTabWidget::pane {
                border: 1px solid #2e3557;
                background: #111424;
                border-radius: 12px;
                padding: 10px;
            }
            QTabBar::tab {
                background: #1a1f3c;
                border: 1px solid #2e3557;
                border-bottom: none;
                border-top-left-radius: 8px;
                border-top-right-radius: 8px;
                padding: 10px 20px;
                font-weight: bold;
                color: #a4b2fc;
            }
            QTabBar::tab:selected {
                background: #252b4d;
                border-color: #4759e5;
                color: #ffffff;
            }
            QTabBar::tab:hover {
                background: #1f2647;
            }
            QWidget#StatsCard1 {
                background: qlineargradient(x1:0, y1:0, x2:1, y2:0, stop:0 #0f172a, stop:1 #1e1b4b);
                border: 1px solid #3b82f6;
                border-radius: 12px;
            }
            QWidget#StatsCard2 {
                background: qlineargradient(x1:0, y1:0, x2:1, y2:0, stop:0 #0f172a, stop:1 #311042);
                border: 1px solid #c084fc;
                border-radius: 12px;
            }
            QWidget#StatsCard3 {
                background: qlineargradient(x1:0, y1:0, x2:1, y2:0, stop:0 #0f172a, stop:1 #064e3b);
                border: 1px solid #10b981;
                border-radius: 12px;
            }
            QProgressBar {
                background-color: #1e293b;
                border: 1px solid #475569;
                border-radius: 6px;
                text-align: right;
                color: transparent;
                height: 12px;
            }
            QProgressBar::chunk {
                background: qlineargradient(x1:0, y1:0, x2:1, y2:0, stop:0 #4f46e5, stop:1 #06b6d4);
                border-radius: 6px;
            }
            """
        )

    def setup_camera(self):
        is_windows = platform.system() == "Windows"
        
        for index in [0, 1, 2]:
            if is_windows:
                cap = cv2.VideoCapture(index, cv2.CAP_DSHOW)
            else:
                cap = cv2.VideoCapture(index)

            if cap.isOpened():
                ok, _ = cap.read()
                if ok:
                    self.cap = cap
                    self.camera_index = index
                    break
                cap.release()

        if not self.cap:
            self.video_label.setText(
                "Camera not opened. Click Start Presenter for manual activation."
            )
            self.set_status("Camera failed or denied. Manual start button is available.")
            return

        self.cap.set(cv2.CAP_PROP_FRAME_WIDTH, 640)
        self.cap.set(cv2.CAP_PROP_FRAME_HEIGHT, 480)

        self.camera_timer = QTimer(self)
        self.camera_timer.timeout.connect(self.update_camera)
        self.camera_timer.start(33)

        self.set_status(f"Camera running. Face detection active.")

    def update_camera(self):
        if not self.cap:
            return

        ok, frame = self.cap.read()
        if not ok:
            self.video_label.setText("Camera frame not received.")
            return

        frame = cv2.flip(frame, 1)
        self.frame_count += 1

        # Real-time Hand Gesture Recognition
        gesture, display_text = self.gesture_detector.process_frame(frame, frame)
        
        if gesture == "Palm Held":
            if self.conv_engine.state == "PRESENTING":
                logger.info("Palm gesture held: Pausing presentation.")
                self.conv_engine.set_state("INTERRUPTED")
                self.conv_engine.interrupt_handler.interrupt(
                    question="Gesture Pause Request",
                    current_slide_index=self.conv_engine.current_slide,
                    slide_title=self.ppt.slides[self.conv_engine.current_slide]["title"] if self.ppt.slides else "Slide",
                    was_auto_presenting=self.conv_engine.auto_presenting
                )
        elif gesture == "Palm Released":
            if self.conv_engine.state in ["QA_SESSION", "INTERRUPTED"]:
                logger.info("Palm gesture released: Resuming presentation.")
                self.conv_engine.interrupt_handler.request_resume()
        elif gesture == "Swipe Left":
            logger.info("Swipe Left gesture detected: Navigating to next slide.")
            self.next_slide_manual()
        elif gesture == "Swipe Right":
            logger.info("Swipe Right gesture detected: Navigating to previous slide.")
            self.prev_slide_manual()

        # Display gesture overlay text
        if display_text != "None":
            cv2.putText(
                frame,
                display_text,
                (20, 70),
                cv2.FONT_HERSHEY_SIMPLEX,
                0.7,
                (255, 100, 100),  # Soft blue/cyan
                2
            )

        gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
        gray = cv2.equalizeHist(gray)

        detections = []

        if self.frontal_cascade is not None:
            faces = self.frontal_cascade.detectMultiScale(
                gray, scaleFactor=1.08, minNeighbors=3, minSize=(40, 40)
            )
            for (x, y, w, h) in faces:
                detections.append(("Visitor", x, y, w, h))

        if not detections and self.profile_cascade is not None:
            profiles = self.profile_cascade.detectMultiScale(
                gray, scaleFactor=1.08, minNeighbors=3, minSize=(40, 40)
            )
            for (x, y, w, h) in profiles:
                detections.append(("Visitor Profile", x, y, w, h))

        if not detections and self.upperbody_cascade is not None:
            bodies = self.upperbody_cascade.detectMultiScale(
                gray, scaleFactor=1.1, minNeighbors=3, minSize=(80, 80)
            )
            for (x, y, w, h) in bodies:
                detections.append(("Visitor Upper Body", x, y, w, h))

        if detections:
            detections.sort(key=lambda item: item[3] * item[4], reverse=True)
            label, x, y, w, h = detections[0]

            self.last_detection_label = label
            self.face_stable_count = min(10, self.face_stable_count + 1)
            self.no_face_frames = 0

            # Calculate relative face coordinates for eye contact pupil tracking
            frame_w = frame.shape[1]
            frame_h = frame.shape[0]
            face_center_x = x + w / 2
            face_center_y = y + h / 2
            
            # Scale coordinates relative to center (rel_x, rel_y from -1.0 to 1.0)
            rel_x = (face_center_x - (frame_w / 2)) / (frame_w / 2)
            rel_y = (face_center_y - (frame_h / 2)) / (frame_h / 2)
            
            # Send coordinates to avatar widget
            self.avatar.set_visitor_position(rel_x, rel_y)

            cv2.rectangle(frame, (x, y), (x + w, y + h), (90, 230, 150), 2)
            cv2.putText(
                frame,
                f"Visitor Detected: {label}",
                (x, max(25, y - 10)),
                cv2.FONT_HERSHEY_SIMPLEX,
                0.7,
                (90, 230, 150),
                2,
            )
        else:
            self.last_detection_label = "None"
            self.face_stable_count = max(0, self.face_stable_count - 1)
            
            # Reset pupil position to center when no visitor is present
            self.avatar.set_visitor_position(0.0, 0.0)
            
            # Count frames with no visitor
            if self.demo_started:
                self.no_face_frames += 1
                if self.no_face_frames > 150: # ~5 seconds
                    self.demo_started = False
                    self.no_face_frames = 0
                    self.analytics.log_session_end()
                    self.conv_engine.handle_visitor_leaving()

            cv2.putText(
                frame,
                "No visitor detected - align face in camera",
                (20, 35),
                cv2.FONT_HERSHEY_SIMPLEX,
                0.65,
                (100, 150, 255),
                2,
            )

        cv2.putText(
            frame,
            f"Stability: {self.face_stable_count}/4 | Mode: {self.last_detection_label}",
            (20, frame.shape[0] - 20),
            cv2.FONT_HERSHEY_SIMPLEX,
            0.6,
            (255, 255, 255),
            2,
        )



        # Trigger automatically when face is stable
        if self.face_stable_count >= 4 and not self.demo_started:
            self.demo_started = True
            self.avatar.pop_out()
            self.analytics.log_visitor()
            self.conv_engine.start_interaction()

        # Display camera stream in label
        rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
        h, w, ch = rgb.shape
        image = QImage(rgb.data, w, h, ch * w, QImage.Format_RGB888)
        pixmap = QPixmap.fromImage(image).scaled(
            self.video_label.width(),
            self.video_label.height(),
            Qt.KeepAspectRatio,
            Qt.SmoothTransformation,
        )
        self.video_label.setPixmap(pixmap)

    # ─── Conversation Engine Callbacks ─────────────────────────────────────────

    def on_state_changed(self, state: str):
        self.avatar.set_state(state.lower(), self.conv_engine.current_slide)
        self.set_status(f"Mode: {state}")
        
        # Update Start Presenter button text dynamically based on interrupted state
        if hasattr(self, "start_btn"):
            if self.conv_engine.interrupt_handler.get_saved_state() is not None:
                self.start_btn.setText("▶️ Resume Presentation")
            else:
                self.start_btn.setText("🚀 Start Presenter")

    def on_text_to_speak(self, text: str, token: str, language: str):
        # Feed text to avatar for phoneme-synced lip movement
        if hasattr(self.avatar, 'set_speaking_text'):
            self.avatar.set_speaking_text(text)
        self.tts.speak(text, token, language)

    def on_speaking_state_changed(self, is_speaking: bool):
        # 1. Lip sync update
        self.avatar.set_speaking(is_speaking)
        
        # 2. Disable voice recognition while agent is speaking to prevent feedback loops
        if self.voice_listener and self.mic_toggle.isChecked():
            self.voice_listener.set_enabled(not is_speaking)

    # ─── UI & Listener Actions ──────────────────────────────────────────────────

    def force_start_demo(self):
        if not self.demo_started:
            self.demo_started = True
            self.avatar.pop_out()
            self.analytics.log_visitor()
            self.conv_engine.start_interaction()
        else:
            self.set_status("Agent is already active.")

    def on_start_or_resume_clicked(self):
        if not self.demo_started:
            self.force_start_demo()
            return
            
        # If active, check if there is an interrupted state to resume
        if self.conv_engine.interrupt_handler.get_saved_state() is not None:
            self.conv_engine.interrupt_handler.request_resume()
        else:
            # Otherwise restart presentation from slide 1
            self.conv_engine.start_presentation()





    def upload_presentation(self):
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            "Open Presentation",
            "",
            "PowerPoint Files (*.pptx)"
        )
        if file_path:
            logger.info(f"Uploading new presentation: {file_path}")
            
            # Stop active speaking and reset engine state to prevent slide index alignment issues
            self.tts.stop_speaking()
            self.conv_engine.reset_engine()
            
            self.ppt.load_presentation(file_path)
            if not self.ppt.slides:
                self.set_status("Failed to load custom presentation.")
                return
            
            # Reset slide UI views
            self.show_slide(0)
            self.set_status(f"Custom presentation loaded: {Path(file_path).name}")
            
            # If the presenter was active, start presentation immediately with the new deck
            if self.demo_started:
                self.conv_engine.start_presentation()

    def select_language(self, code: str):
        self.conv_engine.choose_language(code)
        self.show_slide(self.conv_engine.current_slide)
        
        # Sync checked states of language buttons
        if hasattr(self, 'lang_buttons'):
            for btn, btn_code in self.lang_buttons:
                btn.setChecked(btn_code == code)

    def show_slide(self, index: int):
        if not self.ppt.slides:
            return
        index = max(0, min(index, len(self.ppt.slides) - 1))
        slide = self.ppt.slides[index]
        self.analytics.log_slide_view(slide["title"])
        
        # Translate slide content to the selected language
        lang = self.lm.get_language()
        trans_title, trans_text = self.llm.translate_slide(slide["title"], slide["text"], lang)
        
        self.slide_title.setText(f"Slide {slide['number']}: {trans_title}")
        self.slide_text.setPlainText(trans_text)

    def repeat_slide(self):
        if not self.demo_started:
            self.force_start_demo()
        self.conv_engine.explain_current_slide()

    def next_slide_manual(self):
        if not self.demo_started:
            self.force_start_demo()
        
        # Go to next slide
        if self.conv_engine.current_slide < len(self.ppt.slides) - 1:
            self.conv_engine.current_slide += 1
            self.show_slide(self.conv_engine.current_slide)
            self.conv_engine.explain_current_slide()

    def prev_slide_manual(self):
        if not self.demo_started:
            self.force_start_demo()
        
        # Go to previous slide
        if self.conv_engine.current_slide > 0:
            self.conv_engine.current_slide -= 1
            self.show_slide(self.conv_engine.current_slide)
            self.conv_engine.explain_current_slide()

    def restart_presentation(self):
        if not self.demo_started:
            self.force_start_demo()
        self.conv_engine.start_presentation()

    def handle_typed_question(self):
        question = self.input_box.text().strip()
        self.input_box.clear()
        if question:
            self.add_chat("You", question)
            self.analytics.log_question(question, self.lm.get_language())
            self.conv_engine.handle_question(question)

    def toggle_mic(self):
        if not self.voice_listener:
            self.mic_toggle.setChecked(False)
            return

        active = self.mic_toggle.isChecked()
        if active:
            self.voice_listener.set_enabled(True)
            self.mic_toggle.setText("🎙️ Mic: ON")
        else:
            self.voice_listener.set_enabled(False)
            self.mic_toggle.setText("🎙️ Mic: OFF")

    def on_partial_voice_text(self, text: str):
        # Notify engine user is speaking to immediately pause presentation speech
        self.conv_engine.handle_partial_speech(text)

    def on_voice_text(self, text: str):
        # Full utterance received, print to chat and submit to dialog engine
        self.add_chat("You (Voice)", text)
        self.analytics.log_question(text, self.lm.get_language())
        self.conv_engine.handle_voice_question(text)

    def add_chat(self, sender: str, message: str):
        safe_sender = html.escape(sender)
        safe_message = html.escape(message)
        self.chat.append(f"<b>{safe_sender}:</b> {safe_message}<br>")

    def set_status(self, message: str):
        self.status.setText(f"Status: {message}")

    def closeEvent(self, event):
        try:
            if self.cap:
                self.cap.release()
        except Exception:
            pass

        try:
            if self.voice_listener:
                self.voice_listener.stop()
        except Exception:
            pass

        try:
            self.tts.stop()
        except Exception:
            pass

        event.accept()


def main():
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec())


if __name__ == "__main__":
    main()